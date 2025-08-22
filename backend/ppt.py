import os
import re
import json
import pathlib
from typing import List, Dict, Optional

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Import LLM helpers (optional)
try:
    from llm import get_gemini_client, get_google_search_tool, generate_course_content, system_prompt
except Exception:
    get_gemini_client = None
    get_google_search_tool = None
    generate_course_content = None
    system_prompt = None


def read_all_text_files() -> Dict[str, str]:
    """Read all .txt files from "Inputs and Outputs" and "copilot/Inputs and Outputs".
    Keys are filenames (or prefixed with 'copilot/').
    """
    text_files_content: Dict[str, str] = {}

    home_dir_path = pathlib.Path("Inputs and Outputs")
    if home_dir_path.exists():
        for txt_file in home_dir_path.glob("*.txt"):
            try:
                text_files_content[txt_file.name] = txt_file.read_text(encoding="utf-8")
            except Exception:
                try:
                    text_files_content[txt_file.name] = txt_file.read_text(errors="ignore")
                except Exception:
                    pass

    copilot_dir_path = pathlib.Path("copilot/Inputs and Outputs")
    if copilot_dir_path.exists():
        for txt_file in copilot_dir_path.glob("*.txt"):
            key = f"copilot/{txt_file.name}"
            try:
                text_files_content[key] = txt_file.read_text(encoding="utf-8")
            except Exception:
                try:
                    text_files_content[key] = txt_file.read_text(errors="ignore")
                except Exception:
                    pass

    return text_files_content


def extract_course_name_from_content(content: str) -> str:
    patterns = [
        r'# Course Name:\s*([^\n]+)',
        r'# ([^:\n]+):\s*([^:\n]+)',
        r'# ([^:\n]+)',
        r'## Course Outline:\s*([^\n]+)',
        r'Course Title:\s*([^\n]+)',
        r'Course:\s*([^\n]+)',
        r'Title:\s*([^\n]+)',
    ]
    for pattern in patterns:
        m = re.search(pattern, content, re.IGNORECASE | re.MULTILINE)
        if m:
            course_name = (
                f"{m.group(1).strip()} {m.group(2).strip()}" if len(m.groups()) == 2 else m.group(1).strip()
            )
            course_name = re.sub(r'\s+', ' ', course_name)
            course_name = re.sub(r'[^\w\s-]', '', course_name).strip()
            words = course_name.lower().split()
            key_words = [w for w in words if w not in [
                'advanced', 'introduction', 'to', 'and', 'the', 'of', 'in', 'course', 'fundamentals', 'basics']]
            if key_words:
                subject_name = '_'.join(key_words[:3])
                return f"{subject_name}_course"
    return "course_material"


def _parse_weeks_simple(content: str) -> List[Dict]:
    """Parse '# Week N: ...' blocks with their content, preserving order."""
    weeks: List[Dict] = []
    for m in re.finditer(r"^#\s*Week\s*(\d+)\s*:(.*)$", content, re.MULTILINE | re.IGNORECASE):
        num = int(m.group(1))
        title_line = m.group(0)
        start = m.end()
        next_m = re.search(r"^#\s*Week\s*\d+\s*:.*$", content[start:], re.MULTILINE | re.IGNORECASE)
        end = start + next_m.start() if next_m else len(content)
        block = content[start:end].strip()
        weeks.append({
            'number': num,
            'title': re.sub(r'^#\s*', '', title_line).strip(),
            'content': block,
        })
    weeks.sort(key=lambda x: x['number'])
    if not weeks:
        # Fallback: whole content as Week 1 if nothing found
        weeks = [{
            'number': 1,
            'title': 'Week 1: Course Content',
            'content': content.strip(),
        }]
    return weeks


def sanitize_filename(name: str) -> str:
    name = name.strip().replace('\n', ' ')
    name = re.sub(r"[\\/:*?\"<>|]", "_", name)
    name = re.sub(r"\s+", " ", name)
    return name or "course_material"


def extract_title_from_planner(text: str) -> Optional[str]:
    m = re.search(r"^#\s*Course Name:\s*(.+)$", text, re.MULTILINE)
    if m:
        return m.group(1).strip()
    m = re.search(r"^#\s+(.+)$", text, re.MULTILINE)
    if m:
        return m.group(1).strip()
    m = re.search(r"\*\*Course Name\*\*\s*[:\-]?\s*(.+)$", text, re.IGNORECASE | re.MULTILINE)
    if m:
        return m.group(1).strip()
    return None


def get_subject_from_config() -> Optional[str]:
    try:
        with open('user_config.json', 'r', encoding='utf-8') as f:
            data = json.load(f)
        raw = data.get('subject') or data.get('course_subject') or data.get('course_name')
        if raw and isinstance(raw, str):
            cleaned = re.sub(r'[^\w\s-]', '', raw).strip().lower()
            cleaned = re.sub(r'\s+', ' ', cleaned)
            parts = [p for p in cleaned.split(' ') if p]
            if parts:
                return '_'.join(parts)
    except Exception:
        pass
    return None


# ---------- PPT helpers ----------

def _new_presentation() -> Presentation:
    prs = Presentation()
    # Slightly nicer default text sizes
    return prs


def _add_title_slide(prs: Presentation, course_title: str, week_title: str):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = course_title
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x14, 0x37, 0x66)  # deep blue

    subtitle.text = week_title
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x44, 0x88, 0xCC)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, heading: str, bullets: List[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = heading
    p = body.text_frame.paragraphs[0]
    # Reset first paragraph (it exists by default)
    p.text = bullets[0] if bullets else ""
    p.level = 0
    p.font.size = Pt(18)

    for b in bullets[1:]:
        rp = body.text_frame.add_paragraph()
        rp.text = b
        rp.level = 0
        rp.font.size = Pt(18)


def _chunk_lines(lines: List[str], max_lines: int = 8) -> List[List[str]]:
    chunk: List[List[str]] = []
    current: List[str] = []
    for ln in lines:
        if not ln:
            continue
        current.append(ln)
        if len(current) >= max_lines:
            chunk.append(current)
            current = []
    if current:
        chunk.append(current)
    return chunk


def _text_to_bullets(text: str) -> List[str]:
    """Normalize a paragraph into bullet-sized lines (split sentences)."""
    lines = []
    # Split on explicit bullets first
    for raw in text.split('\n'):
        s = raw.strip()
        if not s:
            continue
        s = re.sub(r'^[-*]\s*', '', s)
        # Further split by sentences if long
        parts = re.split(r'(?<=[.!?])\s+', s)
        for part in parts:
            part = part.strip()
            if part:
                lines.append(part)
    return lines


def _build_week_slide_outline_llm(course_title: str, week_title: str, week_content: str, planner_text: str) -> Optional[List[Dict[str, List[str]]]]:
    """Use LLM to create a clean slide outline for a week.
    Returns a list of dicts: {"title": str, "bullets": [str, ...]} or None on failure.
    """
    if not (get_gemini_client and get_google_search_tool and generate_course_content and system_prompt):
        return None
    try:
        client = get_gemini_client()
        tool = get_google_search_tool()
        task = (
            "You are creating presentation slides for a course week.\n"
            "Given the course title, week title, planner guidance, and raw week text, produce a JSON array named 'slides'.\n"
            "Each item must be an object with keys: 'title' (string) and 'bullets' (array of 4-8 concise strings).\n"
            "Constraints: No markdown, no numbering prefixes unless essential; keep bullets crisp, presentable, and non-redundant.\n"
            "Prefer grouping into logical sections (Concepts, Example, Case Study, Exercise, Tips) if relevant.\n"
            "Output ONLY the JSON, no prose."
        )
        course_content = (
            f"COURSE TITLE: {course_title}\n" 
            f"WEEK TITLE: {week_title}\n\n"
            f"PLANNER INPUT:\n{planner_text}\n\n"
            f"WEEK RAW TEXT:\n{week_content}"
        )
        resp = generate_course_content(
            client=client,
            teaching_style="",
            duration="",
            difficulty_level="",
            google_search_tool=tool,
            system_prompt=system_prompt,
            filepath=pathlib.Path("Inputs and Outputs/curriculum.pdf"),
            course_content=course_content,
            task=task,
        )
        text = resp.text or ""
        # Extract JSON block if wrapped
        json_str = text.strip()
        # Try to find first '[' to last ']' for array
        if '[' in json_str and ']' in json_str:
            json_str = json_str[json_str.find('['): json_str.rfind(']') + 1]
        slides = json.loads(json_str)
        # Validate structure
        cleaned = []
        for item in slides:
            title = str(item.get('title', '')).strip() or 'Section'
            bullets = [str(b).strip() for b in (item.get('bullets') or []) if str(b).strip()]
            if not bullets:
                continue
            cleaned.append({"title": title, "bullets": bullets[:8]})
        return cleaned or None
    except Exception:
        return None


def build_week_ppt(course_title: str, week_title: str, week_content: str, out_path: str, planner_text: str = ""):
    prs = _new_presentation()

    # Title slide
    _add_title_slide(prs, course_title, week_title)

    # Try LLM-authored slide outline first
    llm_slides = _build_week_slide_outline_llm(course_title, week_title, week_content, planner_text)
    if llm_slides:
        for s in llm_slides:
            title = s.get('title', 'Section')
            bullets = s.get('bullets', [])
            for group in _chunk_lines(bullets, max_lines=8):
                _add_content_slide(prs, title, group)
    else:
        # Fallback: parse content into sections by headings
        current_heading = "Overview"
        bucket_lines: List[str] = []

        def flush_bucket():
            nonlocal bucket_lines, current_heading
            if not bucket_lines:
                return
            bullets = _text_to_bullets('\n'.join(bucket_lines))
            for group in _chunk_lines(bullets, max_lines=8):
                _add_content_slide(prs, current_heading, group)
            bucket_lines = []

        for raw in week_content.split('\n'):
            line = raw.rstrip()
            if not line.strip():
                bucket_lines.append(line)
                continue

            if line.lstrip().startswith('### '):
                flush_bucket()
                current_heading = line.lstrip().replace('### ', '').strip()
                continue
            if line.lstrip().startswith('## '):
                flush_bucket()
                current_heading = line.lstrip().replace('## ', '').strip()
                continue
            if line.lstrip().lower().startswith('# week '):
                continue
            bucket_lines.append(line)

        flush_bucket()

    # Save
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


def main():
    # Gather inputs similar to course_material.py
    text_files = read_all_text_files()

    combined_content = ""
    if 'enhanced_course_content.txt' in text_files:
        combined_content = text_files['enhanced_course_content.txt']
    else:
        for filename, content in text_files.items():
            combined_content += f"\n\n=== CONTENT FROM {filename.upper()} ===\n"
            combined_content += content
            combined_content += f"\n=== END OF {filename.upper()} ===\n"

    # Clean out agent/system-like prompts
    lines = combined_content.split('\n')
    cleaned_lines: List[str] = []
    skip_section = False
    for line in lines:
        line_lower = line.lower().strip()
        if any(k in line_lower for k in [
            'teaching agent', 'system prompt', 'agent prompt', 'instruction:',
            'you are a', 'act as', 'your role', 'generate', 'create a course',
            'llm', 'ai assistant', 'chatgpt', 'copilot'
        ]):
            skip_section = True
            continue
        if skip_section and line.strip() == '':
            continue
        elif line.strip() != '':
            skip_section = False
        if not skip_section:
            cleaned_lines.append(line)
    combined_content = '\n'.join(cleaned_lines)

    # Planner title
    planner_path = pathlib.Path("Inputs and Outputs/planner_agent_instruction.txt")
    planner_text = ""
    if planner_path.exists():
        try:
            planner_text = planner_path.read_text(encoding='utf-8')
        except Exception:
            planner_text = planner_path.read_text(errors='ignore')

    course_title = extract_title_from_planner(planner_text) or extract_course_name_from_content(planner_text or combined_content)
    if not course_title or course_title == 'course_material':
        subject = get_subject_from_config()
        if subject:
            course_title = subject.replace('_', ' ').title()
        else:
            course_title = "Course"

    # Parse weeks
    weeks = _parse_weeks_simple(combined_content)

    # Output dir
    out_dir = os.path.join("Inputs and Outputs", "ppts")
    os.makedirs(out_dir, exist_ok=True)

    # Generate one PPT per week
    for w in weeks:
        week_num = w.get('number', 1)
        week_title = w.get('title', f"Week {week_num}")
        content = w.get('content', '')
        filename = f"{sanitize_filename(course_title)}_Week_{week_num:02}.pptx"
        out_path = os.path.join(out_dir, filename)
        build_week_ppt(course_title, week_title, content, out_path, planner_text=planner_text)
        print(f"Created: {out_path}")


if __name__ == "__main__":
    main()
